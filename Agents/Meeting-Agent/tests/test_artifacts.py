import json
import zipfile
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation

from meeting_agent.analyzers import OfflineContractAnalyzer
from meeting_agent.artifacts import (
    BRANCH_COLORS,
    _atomic_generate,
    _clip_display,
    _connector_points,
    _display_width,
    _wrap_display,
    generate_artifacts,
)
from meeting_agent.session import load_jsonl

ROOT = Path(__file__).resolve().parents[1]


def test_generates_nonblank_mind_map_and_valid_pptx(tmp_path: Path) -> None:
    analysis = OfflineContractAnalyzer().analyze(
        load_jsonl(ROOT / "examples" / "product-planning.jsonl")
    )
    artifacts = generate_artifacts(analysis, tmp_path)

    image = Image.open(artifacts["mind_map_png"])
    assert image.size == (1280, 720)
    assert all(variance > 100 for variance in ImageStat.Stat(image).var)

    assert zipfile.is_zipfile(artifacts["presentation"])
    with zipfile.ZipFile(artifacts["presentation"]) as archive:
        assert "ppt/presentation.xml" in archive.namelist()
    presentation = Presentation(artifacts["presentation"])
    assert len(presentation.slides) == 6
    assert "Azure-Agent-Skills-In-Action.pptx" in (
        presentation.core_properties.comments or ""
    )
    text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert analysis.title in text
    assert "Topic content" not in text
    assert "Decision content" not in text
    assert any(shape.shape_type == 13 for shape in presentation.slides[4].shapes)

    graph = json.loads(artifacts["mind_map_json"].read_text(encoding="utf-8"))
    assert graph["label"] == analysis.title
    svg = artifacts["mind_map_svg"].read_text(encoding="utf-8")
    assert svg.startswith("<svg")
    assert 'viewBox="0 0 1280 720"' in svg
    assert '<rect x="450" y="275" width="380" height="170"' in svg
    branch_count = min(6, len(analysis.mind_map.children))
    assert svg.count("<path ") == branch_count
    for color in BRANCH_COLORS[:branch_count]:
        assert f'stroke="{color}"' in svg
    mermaid = artifacts["mind_map_mermaid"].read_text(encoding="utf-8")
    assert mermaid.startswith("mindmap\n  root((")
    assert analysis.title in mermaid.replace("<br/>", " ")
    for text_artifact in ("analysis", "mind_map_json", "mind_map_mermaid", "mind_map_svg"):
        assert b"\r\n" not in artifacts[text_artifact].read_bytes()


def test_atomic_generation_preserves_existing_file_on_failure(tmp_path: Path) -> None:
    target = tmp_path / "artifact.bin"
    target.write_bytes(b"accepted")

    def fail_after_partial_write(temporary: Path) -> None:
        temporary.write_bytes(b"partial")
        raise RuntimeError("generation failed")

    try:
        _atomic_generate(target, fail_after_partial_write)
    except RuntimeError as error:
        assert str(error) == "generation failed"
    else:
        raise AssertionError("expected generation failure")

    assert target.read_bytes() == b"accepted"
    assert list(tmp_path.glob("*.tmp")) == []


def test_wraps_and_clips_cjk_text_by_display_width() -> None:
    value = "会议总结和下一步计划安排讨论会议时间以及风险评估和发布准备"

    clipped = _clip_display(value, 12)
    lines = _wrap_display(value, width=12, max_lines=3)

    assert clipped.endswith("…")
    assert _display_width(clipped) <= 12
    assert 1 < len(lines) <= 3
    assert all(_display_width(line) <= 12 for line in lines)
    assert lines[-1].endswith("…")


def test_routes_lower_connectors_between_card_rows() -> None:
    points = _connector_points(640, 145, 255, 360)

    assert points == [(640, 145), (640, 345), (255, 345), (255, 360)]
    assert 320 < points[1][1] < 360
