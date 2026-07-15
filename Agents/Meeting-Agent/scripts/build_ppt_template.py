"""Build the editable 16:9 Meeting Agent PowerPoint template."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "src" / "meeting_agent" / "templates" / "meeting-agent-template.pptx"

COLORS = {
    "charcoal": "242424",
    "paper": "F7F4EF",
    "white": "FFFFFF",
    "berry": "B11F4B",
    "teal": "007C83",
    "gold": "E7A93B",
    "mist": "E8EEEE",
    "line": "D9D7D2",
    "muted": "68645F",
}


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(COLORS[name])


def add_box(slide, name: str, text: str, x: float, y: float, w: float, h: float, **style):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(style.get("margin", 0))
    frame.margin_right = Inches(style.get("margin", 0))
    frame.margin_top = Inches(style.get("margin", 0))
    frame.margin_bottom = Inches(style.get("margin", 0))
    frame.vertical_anchor = style.get("valign", MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = style.get("align", PP_ALIGN.LEFT)
    paragraph.font.name = style.get("font", "Aptos")
    paragraph.font.size = Pt(style.get("size", 18))
    paragraph.font.bold = style.get("bold", False)
    paragraph.font.color.rgb = rgb(style.get("color", "charcoal"))
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_round_rect(slide, x: float, y: float, w: float, h: float, fill: str, line: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    return shape


def add_header(slide, section: str, title: str, page: int, *, dark: bool = False):
    color = "white" if dark else "charcoal"
    add_box(
        slide,
        f"MA_SECTION_{page}",
        section.upper(),
        0.65,
        0.38,
        2.8,
        0.3,
        size=10,
        bold=True,
        color="berry",
    )
    add_box(
        slide,
        f"MA_SLIDE_TITLE_{page}",
        title,
        0.65,
        0.78,
        11.8,
        0.65,
        size=30,
        bold=True,
        color=color,
    )
    add_box(
        slide,
        f"MA_PAGE_{page}",
        f"{page:02d}",
        12.05,
        7.05,
        0.6,
        0.2,
        size=9,
        color="muted",
        align=PP_ALIGN.RIGHT,
    )


def build_template() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    cover = presentation.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = rgb("charcoal")
    add_rect(cover, 0, 0, 0.18, 7.5, "berry")
    add_rect(cover, 9.55, 0, 3.78, 7.5, "berry")
    add_rect(cover, 10.25, 0.9, 2.45, 1.45, "teal")
    add_rect(cover, 9.85, 2.65, 2.85, 1.05, "gold")
    add_rect(cover, 10.55, 4.05, 2.15, 2.3, "paper")
    add_box(
        cover,
        "MA_TITLE",
        "Meeting title",
        0.82,
        1.45,
        7.9,
        1.55,
        size=32,
        bold=True,
        color="white",
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_box(
        cover,
        "MA_SUBTITLE",
        "Executive meeting package",
        0.85,
        3.25,
        7.4,
        1.15,
        size=18,
        color="white",
    )
    add_box(
        cover,
        "MA_COVER_LABEL",
        "MEETING AGENT  /  EVIDENCE TO ACTION",
        0.85,
        0.62,
        7.8,
        0.35,
        size=10,
        bold=True,
        color="gold",
    )
    add_box(
        cover,
        "MA_COVER_FOOTER",
        "Generated from validated meeting evidence",
        0.85,
        6.78,
        7.5,
        0.25,
        size=10,
        color="white",
    )

    summary = presentation.slides.add_slide(blank)
    summary.background.fill.solid()
    summary.background.fill.fore_color.rgb = rgb("paper")
    add_header(summary, "01 / Readout", "Executive summary", 2)
    add_round_rect(summary, 0.65, 1.75, 7.55, 4.7, "white", "line")
    add_box(summary, "MA_SUMMARY", "Summary", 1.0, 2.15, 6.85, 3.85, size=20, color="charcoal")
    stats = [
        ("MA_TOPIC_COUNT", "00", "TOPICS", "teal"),
        ("MA_DECISION_COUNT", "00", "DECISIONS", "berry"),
        ("MA_ACTION_COUNT", "00", "ACTIONS", "gold"),
    ]
    for index, (name, value, label, accent) in enumerate(stats):
        y = 1.75 + index * 1.58
        add_round_rect(summary, 8.65, y, 4.0, 1.28, "white", "line")
        add_rect(summary, 8.65, y, 0.12, 1.28, accent)
        add_box(summary, name, value, 9.05, y + 0.18, 1.1, 0.6, size=27, bold=True, color=accent)
        add_box(
            summary,
            f"{name}_LABEL",
            label,
            10.15,
            y + 0.43,
            1.95,
            0.25,
            size=10,
            bold=True,
            color="muted",
        )

    topics = presentation.slides.add_slide(blank)
    topics.background.fill.solid()
    topics.background.fill.fore_color.rgb = rgb("paper")
    add_header(topics, "02 / Themes", "What shaped the discussion", 3)
    accents = ["berry", "teal", "gold"]
    for index, accent in enumerate(accents):
        x = 0.65 + index * 4.18
        add_round_rect(topics, x, 1.78, 3.78, 4.75, "white", "line")
        add_rect(topics, x, 1.78, 3.78, 0.14, accent)
        add_box(
            topics,
            f"MA_TOPIC_INDEX_{index + 1}",
            f"0{index + 1}",
            x + 0.32,
            2.15,
            0.65,
            0.4,
            size=16,
            bold=True,
            color=accent,
        )
        add_box(
            topics,
            f"MA_TOPICS_{index + 1}",
            "Topic content",
            x + 0.32,
            2.75,
            3.12,
            3.25,
            size=17,
            color="charcoal",
        )

    execution = presentation.slides.add_slide(blank)
    execution.background.fill.solid()
    execution.background.fill.fore_color.rgb = rgb("paper")
    add_header(execution, "03 / Execution", "Decisions and action register", 4)
    add_round_rect(execution, 0.65, 1.75, 5.45, 4.8, "white", "line")
    add_rect(execution, 0.65, 1.75, 5.45, 0.6, "berry")
    add_box(
        execution,
        "MA_DECISIONS_LABEL",
        "DECISIONS",
        0.95,
        1.94,
        2.2,
        0.25,
        size=11,
        bold=True,
        color="white",
    )
    add_box(execution, "MA_DECISIONS", "Decision content", 0.98, 2.68, 4.8, 3.35, size=17)
    add_round_rect(execution, 6.45, 1.75, 6.23, 4.8, "white", "line")
    add_rect(execution, 6.45, 1.75, 6.23, 0.6, "teal")
    add_box(
        execution,
        "MA_ACTIONS_LABEL",
        "ACTION REGISTER",
        6.75,
        1.94,
        2.6,
        0.25,
        size=11,
        bold=True,
        color="white",
    )
    add_box(execution, "MA_ACTIONS", "Action content", 6.78, 2.68, 5.56, 3.35, size=16)

    mind_map = presentation.slides.add_slide(blank)
    mind_map.background.fill.solid()
    mind_map.background.fill.fore_color.rgb = rgb("charcoal")
    add_header(mind_map, "04 / Map", "Meeting landscape", 5, dark=True)
    frame = add_round_rect(mind_map, 0.72, 1.62, 11.9, 5.12, "white", "white")
    frame.name = "MA_MIND_MAP_FRAME"
    add_box(
        mind_map,
        "MA_MIND_MAP_NOTE",
        "Rendered from the structured meeting graph",
        0.82,
        6.9,
        4.8,
        0.22,
        size=9,
        color="white",
    )

    questions = presentation.slides.add_slide(blank)
    questions.background.fill.solid()
    questions.background.fill.fore_color.rgb = rgb("paper")
    add_header(questions, "05 / Forward", "Open questions and next conversation", 6)
    add_round_rect(questions, 0.65, 1.75, 7.6, 4.8, "white", "line")
    add_box(questions, "MA_QUESTIONS", "Open questions", 1.0, 2.15, 6.9, 3.7, size=18)
    add_rect(questions, 8.65, 1.75, 4.0, 4.8, "teal")
    add_box(
        questions,
        "MA_NEXT_LABEL",
        "NEXT CONVERSATION",
        9.0,
        2.18,
        3.25,
        0.35,
        size=11,
        bold=True,
        color="white",
    )
    add_box(
        questions,
        "MA_NEXT_STEP",
        "Resolve the open questions and confirm owners.",
        9.0,
        2.9,
        3.15,
        2.4,
        size=22,
        bold=True,
        color="white",
    )
    add_box(
        questions,
        "MA_CLOSE",
        "Evidence first. Human review required.",
        9.0,
        5.75,
        3.15,
        0.4,
        size=11,
        color="white",
    )
    return presentation


def main() -> int:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_template()
    presentation.save(OUTPUT)
    print(f"PPT_TEMPLATE_BUILT path={OUTPUT} slides={len(presentation.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
