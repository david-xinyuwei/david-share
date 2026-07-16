"""Generate renderer-neutral mind-map files and an editable PPTX."""

from __future__ import annotations

import json
import platform
import re
from collections.abc import Callable
from html import escape
from importlib.resources import as_file, files
from pathlib import Path
from unicodedata import east_asian_width
from uuid import uuid4

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import MeetingAnalysis, MindMapNode

BRANCH_COLORS = ("#0078D4", "#00B04F", "#FFAA00", "#7C3AED", "#E74C3C", "#008C95")


def generate_artifacts(analysis: MeetingAnalysis, output_dir: Path) -> dict[str, Path]:
    generated = generate_mind_map_artifacts(analysis, output_dir)
    generated["presentation"] = generate_presentation_artifact(
        analysis,
        generated["mind_map_png"],
        output_dir,
    )
    return generated


def generate_mind_map_artifacts(
    analysis: MeetingAnalysis,
    output_dir: Path,
) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    graph_json = output_dir / "mind-map.json"
    graph_mermaid = output_dir / "mind-map.mmd"
    graph_svg = output_dir / "mind-map.svg"
    graph_png = output_dir / "mind-map.png"
    analysis_json = output_dir / "meeting-analysis.json"

    _atomic_write_text(
        graph_json,
        json.dumps(analysis.mind_map.model_dump(), ensure_ascii=False, indent=2),
    )
    _atomic_write_text(analysis_json, analysis.model_dump_json(indent=2))
    _atomic_write_text(graph_mermaid, mind_map_mermaid(analysis.mind_map))
    _atomic_write_text(graph_svg, _mind_map_svg(analysis.mind_map))
    _atomic_generate(graph_png, lambda temporary: _mind_map_png(analysis.mind_map, temporary))
    return {
        "analysis": analysis_json,
        "mind_map_json": graph_json,
        "mind_map_mermaid": graph_mermaid,
        "mind_map_svg": graph_svg,
        "mind_map_png": graph_png,
    }


def generate_presentation_artifact(
    analysis: MeetingAnalysis,
    mind_map_path: Path,
    output_dir: Path,
) -> Path:
    presentation = output_dir / "meeting-summary.pptx"
    _atomic_generate(
        presentation,
        lambda temporary: _presentation(analysis, mind_map_path, temporary),
    )
    return presentation


def _atomic_write_text(path: Path, text: str) -> None:
    temporary = _temporary_path(path)
    try:
        temporary.write_text(text, encoding="utf-8", newline="\n")
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _atomic_generate(path: Path, writer: Callable[[Path], None]) -> None:
    temporary = _temporary_path(path)
    try:
        writer(temporary)
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _temporary_path(path: Path) -> Path:
    return path.with_name(f".{path.name}.{uuid4().hex}.tmp")


def _mind_map_svg(root: MindMapNode) -> str:
    children = root.children[:6]
    positions = (
        (35, 25, "left"),
        (35, 270, "left"),
        (35, 515, "left"),
        (865, 25, "right"),
        (865, 270, "right"),
        (865, 515, "right"),
    )
    lines = [
        '<svg xmlns="http://www.w3.org/2000/svg" width="1280" '
        'height="720" viewBox="0 0 1280 720">',
        '<rect width="1280" height="720" fill="#F3F4F8"/>',
        "<style>text{font-family:'Segoe UI',Arial,sans-serif;fill:#1A1A3E}"
        ".eyebrow{font-size:12px;font-weight:700;letter-spacing:1px;fill:#0078D4}"
        ".root{font-size:24px;font-weight:700;fill:white}"
        ".node{font-size:17px;font-weight:700}.index{font-size:12px;font-weight:700}"
        ".leaf{font-size:12.5px;fill:#606070}</style>",
        '<text class="eyebrow" x="640" y="28" text-anchor="middle">MEETING LANDSCAPE</text>',
        '<rect x="450" y="275" width="380" height="170" rx="18" fill="#1A1A3E"/>',
        '<rect x="450" y="275" width="380" height="8" rx="4" fill="#0078D4"/>',
    ]
    root_lines = _wrap_display(root.label, width=28, max_lines=3)
    root_start = 350 - (len(root_lines) - 1) * 16
    lines.append('<text class="root" x="640" text-anchor="middle">')
    for line_index, line in enumerate(root_lines):
        lines.append(
            f'<tspan x="640" y="{root_start + line_index * 32}">'
            f"{escape(line)}</tspan>"
        )
    lines.append("</text>")
    lines.append(
        '<text x="640" y="414" text-anchor="middle" '
        'style="font-size:12px;fill:#CCCCDD">EVIDENCE TO ACTION</text>'
    )
    for index, child in enumerate(children):
        x, y, side = positions[index]
        color = BRANCH_COLORS[index]
        card_mid_y = y + 90
        root_x = 450 if side == "left" else 830
        card_x = x + 380 if side == "left" else x
        root_y = 305 + (index % 3) * 55
        control_x = 420 if side == "left" else 860
        lines.extend(
            [
                f'<path d="M{root_x} {root_y} C{control_x} {root_y}, '
                f'{control_x} {card_mid_y}, {card_x} {card_mid_y}" fill="none" '
                f'stroke="{color}" stroke-width="3"/>',
                f'<rect x="{x}" y="{y}" width="380" height="180" rx="14" '
                f'fill="#FFFFFF" stroke="{color}" stroke-width="2"/>',
                f'<rect x="{x}" y="{y}" width="8" height="180" rx="4" fill="{color}"/>',
                f'<text class="index" x="{x + 28}" y="{y + 32}" fill="{color}">'
                f"{index + 1:02d}</text>",
                f'<text class="node" x="{x + 66}" y="{y + 32}">'
                f"{escape(_clip_display(child.label, 32))}</text>",
            ]
        )
        for leaf_index, leaf in enumerate(child.children[:4]):
            leaf_y = y + 68 + leaf_index * 27
            lines.append(
                f'<circle cx="{x + 31}" cy="{leaf_y - 4}" r="3" fill="{color}"/>'
            )
            lines.append(
                f'<text class="leaf" x="{x + 44}" y="{leaf_y}">'
                f"{escape(_clip_display(leaf.label, 48))}</text>"
            )
    lines.append("</svg>")
    return "\n".join(lines)


def mind_map_mermaid(root: MindMapNode) -> str:
    lines = ["mindmap", f"  root(({_mermaid_label(root.label)}))"]
    for child in root.children[:6]:
        lines.append(f"    {_mermaid_label(child.label)}")
        for leaf in child.children[:4]:
            lines.append(f"      {_mermaid_label(leaf.label)}")
    return "\n".join(lines) + "\n"


def _mermaid_label(value: str) -> str:
    compact = " ".join(value.split())
    sanitized = re.sub(r"[()[\]{}\"`]", "", compact) or "Meeting"
    return "<br/>".join(_wrap_all_display(sanitized, width=24))


def _wrap_all_display(value: str, width: int) -> list[str]:
    remaining = value
    lines: list[str] = []
    while remaining:
        if _display_width(remaining) <= width:
            lines.append(remaining)
            break
        used = 0
        cut = 0
        for index, character in enumerate(remaining):
            character_width = 2 if east_asian_width(character) in {"F", "W"} else 1
            if used + character_width > width:
                break
            used += character_width
            cut = index + 1
        segment = remaining[:cut]
        last_space = segment.rfind(" ")
        if last_space > 0:
            cut = last_space
            segment = remaining[:cut]
        lines.append(segment.rstrip())
        remaining = remaining[cut:].lstrip()
    return lines or ["Meeting"]


def _mind_map_png(root: MindMapNode, path: Path) -> None:
    image = Image.new("RGB", (1280, 720), "#F3F4F8")
    draw = ImageDraw.Draw(image)
    title_font = _font(24, bold=True)
    section_font = _font(17, bold=True)
    body_font = _font(12)
    index_font = _font(12, bold=True)
    eyebrow_font = _font(12, bold=True)
    draw.text((640, 22), "MEETING LANDSCAPE", fill="#0078D4", font=eyebrow_font, anchor="mm")
    draw.rounded_rectangle((450, 275, 830, 445), radius=18, fill="#1A1A3E")
    draw.rounded_rectangle((450, 275, 830, 283), radius=4, fill="#0078D4")
    root_text = "\n".join(_wrap_display(root.label, width=28, max_lines=3))
    draw.multiline_text(
        (640, 352),
        root_text,
        fill="white",
        font=title_font,
        anchor="mm",
        align="center",
        spacing=5,
    )
    draw.text((640, 416), "EVIDENCE TO ACTION", fill="#CCCCDD", font=eyebrow_font, anchor="mm")
    positions = (
        (35, 25, "left"),
        (35, 270, "left"),
        (35, 515, "left"),
        (865, 25, "right"),
        (865, 270, "right"),
        (865, 515, "right"),
    )
    for index, child in enumerate(root.children[:6]):
        left, top, side = positions[index]
        right = left + 380
        bottom = top + 180
        color = BRANCH_COLORS[index]
        root_x = 450 if side == "left" else 830
        card_x = right if side == "left" else left
        root_y = 305 + (index % 3) * 55
        card_y = top + 90
        bend_x = 420 if side == "left" else 860
        draw.line(
            (root_x, root_y, bend_x, root_y, bend_x, card_y, card_x, card_y),
            fill=color,
            width=3,
        )
        draw.rounded_rectangle(
            (left, top, right, bottom),
            radius=14,
            fill="#FFFFFF",
            outline=color,
            width=2,
        )
        draw.rounded_rectangle((left, top, left + 8, bottom), radius=4, fill=color)
        draw.text(
            (left + 28, top + 20),
            f"{index + 1:02d}",
            fill=color,
            font=index_font,
        )
        draw.text(
            (left + 66, top + 18),
            _clip_display(child.label, 32),
            fill="#1A1A3E",
            font=section_font,
        )
        for leaf_index, leaf in enumerate(child.children[:4]):
            leaf_y = top + 67 + leaf_index * 27
            draw.ellipse((left + 28, leaf_y, left + 34, leaf_y + 6), fill=color)
            draw.text(
                (left + 44, leaf_y - 5),
                _clip_display(leaf.label, 48),
                fill="#606070",
                font=body_font,
            )
    image.save(path, format="PNG")


def _font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    windows_name = "msyhbd.ttc" if bold else "msyh.ttc"
    noto_name = "NotoSansCJK-Bold.ttc" if bold else "NotoSansCJK-Regular.ttc"
    dejavu_name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    candidates = [
        Path("/usr/share/fonts/opentype/noto") / noto_name,
        Path("/usr/share/fonts/truetype/dejavu") / dejavu_name,
        Path(dejavu_name),
    ]
    if platform.system() == "Windows":
        candidates.insert(0, Path("C:/Windows/Fonts") / windows_name)
    for candidate in candidates:
        try:
            return ImageFont.truetype(str(candidate), size)
        except OSError:
            continue
    return ImageFont.load_default(size=size)


def _presentation(analysis: MeetingAnalysis, mind_map_path: Path, path: Path) -> None:
    template = files("meeting_agent").joinpath("templates/meeting-agent-template.pptx")
    with as_file(template) as template_path:
        presentation = Presentation(template_path)

    _set_shape_text(
        presentation,
        "MA_TITLE",
        _clip_display(analysis.title, 70),
        size=32,
        bold=True,
        color="FFFFFF",
    )
    _set_shape_text(
        presentation,
        "MA_SUBTITLE",
        _clip_display(analysis.summary, 260),
        size=18,
        color="FFFFFF",
    )
    _set_shape_text(presentation, "MA_SUMMARY", analysis.summary, size=16)
    _set_shape_text(
        presentation,
        "MA_TOPIC_COUNT",
        f"{len(analysis.topics):02d}",
        size=27,
        bold=True,
        color="00B04F",
    )
    _set_shape_text(
        presentation,
        "MA_DECISION_COUNT",
        f"{len(analysis.decisions):02d}",
        size=27,
        bold=True,
        color="0078D4",
    )
    _set_shape_text(
        presentation,
        "MA_ACTION_COUNT",
        f"{len(analysis.action_items):02d}",
        size=27,
        bold=True,
        color="FFAA00",
    )

    for index in range(1, 7):
        if index <= len(analysis.topics):
            _set_shape_text(
                presentation,
                f"MA_TOPICS_{index}",
                _clip_display(analysis.topics[index - 1], 110),
                size=16,
                bold=True,
            )
        else:
            _remove_named_shapes(
                presentation,
                (
                    f"MA_TOPIC_CARD_{index}",
                    f"MA_TOPIC_STRIPE_{index}",
                    f"MA_TOPIC_INDEX_{index}",
                    f"MA_TOPICS_{index}",
                ),
            )
    _set_numbered_list(
        presentation,
        "MA_DECISIONS",
        analysis.decisions or ["No decision recorded"],
        size=16,
        limit=5,
    )
    _set_action_register(
        presentation,
        "MA_ACTIONS",
        analysis.action_items,
        limit=5,
    )
    _add_contained_picture(presentation, "MA_MIND_MAP_FRAME", mind_map_path)
    questions = analysis.open_questions or ["No open question recorded"]
    _set_numbered_list(presentation, "MA_QUESTIONS", questions, size=18, limit=7)
    _set_shape_text(
        presentation,
        "MA_NEXT_STEP",
        _clip_display(
            analysis.open_questions[0]
            if analysis.open_questions
            else "Confirm the action owners and publish the reviewed meeting package.",
            150,
        ),
        size=22,
        bold=True,
        color="FFFFFF",
    )
    presentation.save(path)


def _find_shape(presentation: Presentation, name: str):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.name == name:
                return shape
    raise ValueError(f"PowerPoint template shape is missing: {name}")


def _remove_named_shapes(presentation: Presentation, names: tuple[str, ...]) -> None:
    for name in names:
        shape = _find_shape(presentation, name)
        shape._element.getparent().remove(shape._element)


def _add_contained_picture(
    presentation: Presentation,
    frame_name: str,
    image_path: Path,
) -> None:
    frame = _find_shape(presentation, frame_name)
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    margin = Inches(0.14)
    available_width = frame.width - 2 * margin
    available_height = frame.height - 2 * margin
    available_ratio = available_width / available_height
    if image_ratio >= available_ratio:
        width = available_width
        height = int(width / image_ratio)
    else:
        height = available_height
        width = int(height * image_ratio)
    left = frame.left + (frame.width - width) // 2
    top = frame.top + (frame.height - height) // 2
    frame._element.getparent().remove(frame._element)
    presentation.slides[4].shapes.add_picture(
        str(image_path),
        left,
        top,
        width=width,
        height=height,
    )


def _set_shape_text(
    presentation: Presentation,
    name: str,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color: str = "242424",
) -> None:
    shape = _find_shape(presentation, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Segoe UI"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)


def _set_numbered_list(
    presentation: Presentation,
    name: str,
    items: list[str],
    *,
    size: int,
    limit: int,
) -> None:
    shape = _find_shape(presentation, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items[:limit]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"{index + 1:02d}  {_clip_display(item, 180)}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor.from_string("242424")
        paragraph.space_after = Pt(12)


def _set_action_register(
    presentation: Presentation,
    name: str,
    items: list,
    *,
    limit: int,
) -> None:
    shape = _find_shape(presentation, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    if not items:
        paragraph = frame.paragraphs[0]
        paragraph.text = "No action item was identified in the supplied evidence."
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor.from_string("606070")
        return
    for index, item in enumerate(items[:limit]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        number = paragraph.add_run()
        number.text = f"{index + 1:02d}  "
        number.font.name = "Segoe UI"
        number.font.size = Pt(11)
        number.font.bold = True
        number.font.color.rgb = RGBColor.from_string("00B04F")
        content = paragraph.add_run()
        content.text = _clip_display(item.description, 180)
        content.font.name = "Segoe UI"
        content.font.size = Pt(11)
        content.font.color.rgb = RGBColor.from_string("242424")
        metadata = " · ".join(
            value
            for value in (
                f"Owner: {item.owner}" if item.owner else "Owner: unassigned",
                f"Due: {item.due}" if item.due else "Due: not specified",
            )
        )
        details = paragraph.add_run()
        details.text = f"  /  {metadata}"
        details.font.name = "Segoe UI"
        details.font.size = Pt(10)
        details.font.bold = True
        details.font.color.rgb = RGBColor.from_string("606070")
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(7)


def _display_width(value: str) -> int:
    return sum(2 if east_asian_width(character) in {"F", "W"} else 1 for character in value)


def _connector_points(
    source_x: int,
    source_y: int,
    target_x: int,
    target_y: int,
) -> list[tuple[int, int]]:
    branch_y = target_y - 15
    return [
        (source_x, source_y),
        (source_x, branch_y),
        (target_x, branch_y),
        (target_x, target_y),
    ]


def _clip_display(value: str, width: int) -> str:
    compact = " ".join(value.split())
    if _display_width(compact) <= width:
        return compact
    output: list[str] = []
    used = 0
    for character in compact:
        character_width = 2 if east_asian_width(character) in {"F", "W"} else 1
        if used + character_width > width - 1:
            break
        output.append(character)
        used += character_width
    return "".join(output).rstrip() + "…"


def _wrap_display(value: str, width: int, max_lines: int) -> list[str]:
    remaining = " ".join(value.split())
    lines: list[str] = []
    while remaining and len(lines) < max_lines:
        if _display_width(remaining) <= width:
            lines.append(remaining)
            remaining = ""
            break
        used = 0
        cut = 0
        for index, character in enumerate(remaining):
            character_width = 2 if east_asian_width(character) in {"F", "W"} else 1
            if used + character_width > width:
                break
            used += character_width
            cut = index + 1
        segment = remaining[:cut]
        last_space = segment.rfind(" ")
        if last_space > 0:
            cut = last_space
            segment = remaining[:cut]
        lines.append(segment.rstrip())
        remaining = remaining[cut:].lstrip()
    if remaining:
        lines[-1] = _clip_display(f"{lines[-1]} {remaining}", width)
    return lines
