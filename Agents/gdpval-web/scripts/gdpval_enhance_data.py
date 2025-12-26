"""
GDPVAL 数据增强脚本
下载附件并解析内容，合并到 prompt 中生成增强数据集

附件统计:
- .xlsx: 85 → pandas 转 Markdown 表格
- .pdf: 77 → PyPDF2 提取文字
- .docx: 66 → python-docx 提取文字
- .png: 10 → Azure CV OCR 提取文字
- .jpg: 3 → Azure CV OCR 提取文字
- .zip: 3 → 解压后递归处理
- .txt: 1 → 直接读取
- .pptx: 1 → python-pptx 提取文字
- .wav: 8 → ⏭️ 跳过（音频）
- .mp3: 2 → ⏭️ 跳过（音频）
- .mp4: 2 → ⏭️ 跳过（视频）
- .step: 2 → ⏭️ 跳过（CAD）
- .pages: 1 → ⏭️ 跳过（Apple Pages）
- .psd: 1 → ⏭️ 跳过（Photoshop）
"""

import json
import os
import requests
import hashlib
import base64
from pathlib import Path
from datetime import datetime
import pandas as pd
from io import BytesIO
import zipfile
import tempfile

# 可选依赖
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False
    print("⚠️  PyPDF2 未安装，PDF 解析将使用备用方案")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️  python-docx 未安装，Word 解析将跳过")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️  python-pptx 未安装，PowerPoint 解析将跳过")

# ============================================================
# 配置
# ============================================================

CACHE_DIR = Path(__file__).parent / "gdpval_cache"
CACHE_DIR.mkdir(exist_ok=True)

MAX_CONTENT_LENGTH = 50000  # 每个附件最大字符数
MAX_TOTAL_CONTENT = 100000  # 总附件内容最大字符数

# Azure OpenAI 配置 (用于图片分析)
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT", "https://your-aoai-endpoint.openai.azure.com/")
AZURE_OPENAI_KEY = os.environ.get("AZURE_OPENAI_KEY", "YOUR_AZURE_OPENAI_KEY")
AZURE_OPENAI_DEPLOYMENT = "gpt-5.2-chat"  # GPT-5.2 视觉模型

# 支持的文件类型
SUPPORTED_EXTENSIONS = {'.xlsx', '.pdf', '.docx', '.txt', '.pptx', '.zip', '.png', '.jpg', '.jpeg'}
SKIP_EXTENSIONS = {'.wav', '.mp3', '.mp4', '.step', '.psd', '.pages', '.gif'}  # 音视频和特殊格式跳过
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg'}  # 图片用 GPT-5.2 视觉分析


# ============================================================
# 文件下载
# ============================================================

def get_cache_path(url: str) -> Path:
    """生成缓存文件路径"""
    url_hash = hashlib.md5(url.encode()).hexdigest()[:12]
    filename = url.split('/')[-1]
    return CACHE_DIR / f"{url_hash}_{filename}"


def download_file(url: str) -> bytes | None:
    """下载文件，支持缓存"""
    cache_path = get_cache_path(url)
    
    # 检查缓存
    if cache_path.exists():
        print(f"  ✅ 缓存命中: {cache_path.name}")
        return cache_path.read_bytes()
    
    # 下载
    try:
        print(f"  ⬇️  下载: {url.split('/')[-1]}")
        response = requests.get(url, timeout=60)
        response.raise_for_status()
        
        # 保存缓存
        cache_path.write_bytes(response.content)
        return response.content
    except Exception as e:
        print(f"  ❌ 下载失败: {e}")
        return None


# ============================================================
# 文件解析
# ============================================================

def parse_excel(content: bytes, filename: str) -> str:
    """解析 Excel 文件"""
    try:
        # 读取所有 sheet
        xlsx = pd.ExcelFile(BytesIO(content))
        result_parts = []
        
        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            
            if df.empty:
                continue
                
            result_parts.append(f"\n=== Sheet: {sheet_name} ===")
            
            # 转换为 Markdown 表格格式
            result_parts.append(df.to_markdown(index=False))
        
        result = '\n'.join(result_parts)
        
        # 限制长度
        if len(result) > MAX_CONTENT_LENGTH:
            result = result[:MAX_CONTENT_LENGTH] + f"\n... [内容截断，共 {len(result)} 字符]"
        
        return result
        
    except Exception as e:
        return f"[Excel 解析错误: {e}]"


def parse_pdf(content: bytes, filename: str) -> str:
    """解析 PDF 文件 - 优先文本提取，失败则用 GPT-5.2 视觉分析"""
    if not HAS_PYPDF2:
        return "[PDF 解析需要 PyPDF2: pip install PyPDF2]"
    
    try:
        reader = PyPDF2.PdfReader(BytesIO(content))
        text_parts = []
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_parts.append(f"\n--- Page {i+1} ---\n{text}")
        
        result = '\n'.join(text_parts)
        
        if len(result) > MAX_CONTENT_LENGTH:
            result = result[:MAX_CONTENT_LENGTH] + f"\n... [内容截断，共 {len(result)} 字符]"
        
        # 如果文本提取成功，返回结果
        if result.strip():
            return result
        
        # 文本提取失败，尝试用 GPT-5.2 视觉分析 PDF 页面
        print(f"    📄 PDF 无文本层，尝试 GPT-5.2 视觉分析...")
        return parse_pdf_with_vision(content, filename)
        
    except Exception as e:
        return f"[PDF 解析错误: {e}]"


def parse_pdf_with_vision(content: bytes, filename: str) -> str:
    """使用 GPT-5.2 视觉模型分析扫描件 PDF"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "[PDF 视觉分析需要 PyMuPDF: pip install pymupdf]"
    
    try:
        doc = fitz.open(stream=content, filetype="pdf")
        text_parts = []
        
        # 最多分析前 5 页（避免成本过高）
        max_pages = min(len(doc), 5)
        
        for i in range(max_pages):
            page = doc[i]
            
            # 将 PDF 页面渲染为图片
            mat = fitz.Matrix(2, 2)  # 2x 缩放提高清晰度
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            
            # 用 GPT-5.2 分析这一页
            analysis = _analyze_image_with_gpt52(img_bytes, f"{filename}_page{i+1}.png")
            
            if analysis and not analysis.startswith("["):
                text_parts.append(f"\n--- Page {i+1} (Vision) ---\n{analysis}")
        
        doc.close()
        
        result = '\n'.join(text_parts)
        
        if len(result) > MAX_CONTENT_LENGTH:
            result = result[:MAX_CONTENT_LENGTH] + f"\n... [内容截断]"
        
        return result if result.strip() else "[PDF 视觉分析无结果]"
        
    except Exception as e:
        return f"[PDF 视觉分析错误: {e}]"


def _analyze_image_with_gpt52(image_bytes: bytes, filename: str) -> str:
    """调用 GPT-5.2 分析单张图片（内部函数）"""
    try:
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
        
        headers = {
            'api-key': AZURE_OPENAI_KEY,
            'Content-Type': 'application/json'
        }
        
        payload = {
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Extract ALL text and data from this document image. Include tables, headers, numbers, and any visible content. Format tables as markdown. Be thorough."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_base64}"
                            }
                        }
                    ]
                }
            ],
            "max_completion_tokens": 4096
        }
        
        api_url = f"{AZURE_OPENAI_ENDPOINT.rstrip('/')}/openai/deployments/{AZURE_OPENAI_DEPLOYMENT}/chat/completions?api-version=2024-10-21"
        
        response = requests.post(api_url, headers=headers, json=payload, timeout=120)
        
        if response.status_code == 200:
            result = response.json()
            return result.get('choices', [{}])[0].get('message', {}).get('content', '')
        else:
            return f"[Vision API 错误: HTTP {response.status_code}]"
        
    except Exception as e:
        return f"[Vision 分析错误: {e}]"


def parse_docx(content: bytes, filename: str) -> str:
    """解析 Word 文档"""
    if not HAS_DOCX:
        return "[Word 解析需要 python-docx: pip install python-docx]"
    
    try:
        doc = Document(BytesIO(content))
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 也提取表格
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                table_rows.append(row_text)
            if table_rows:
                text_parts.append('\n[表格]\n' + '\n'.join(table_rows))
        
        result = '\n'.join(text_parts)
        
        if len(result) > MAX_CONTENT_LENGTH:
            result = result[:MAX_CONTENT_LENGTH] + f"\n... [内容截断，共 {len(result)} 字符]"
        
        return result if result.strip() else "[Word 无可提取文本]"
        
    except Exception as e:
        return f"[Word 解析错误: {e}]"


def parse_pptx(content: bytes, filename: str) -> str:
    """解析 PowerPoint 文件"""
    if not HAS_PPTX:
        return "[PowerPoint 解析需要 python-pptx: pip install python-pptx]"
    
    try:
        prs = Presentation(BytesIO(content))
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"\n--- Slide {i+1} ---\n" + '\n'.join(slide_texts))
        
        result = '\n'.join(text_parts)
        
        if len(result) > MAX_CONTENT_LENGTH:
            result = result[:MAX_CONTENT_LENGTH] + f"\n... [内容截断]"
        
        return result if result.strip() else "[PowerPoint 无可提取文本]"
        
    except Exception as e:
        return f"[PowerPoint 解析错误: {e}]"


def parse_txt(content: bytes, filename: str) -> str:
    """解析文本文件"""
    try:
        # 尝试多种编码
        for encoding in ['utf-8', 'gbk', 'latin-1']:
            try:
                text = content.decode(encoding)
                if len(text) > MAX_CONTENT_LENGTH:
                    text = text[:MAX_CONTENT_LENGTH] + f"\n... [内容截断]"
                return text
            except:
                continue
        return "[文本编码无法识别]"
    except Exception as e:
        return f"[文本解析错误: {e}]"


def parse_image(content: bytes, filename: str) -> str:
    """使用 GPT-5.2 视觉模型分析图片内容"""
    try:
        # 将图片转为 base64
        image_base64 = base64.b64encode(content).decode('utf-8')
        
        # 根据扩展名确定 MIME 类型
        ext = Path(filename).suffix.lower()
        mime_type = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg'
        }.get(ext, 'image/png')
        
        # 调用 Azure OpenAI GPT-5.2 视觉 API
        headers = {
            'api-key': AZURE_OPENAI_KEY,
            'Content-Type': 'application/json'
        }
        
        payload = {
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "Please analyze this image and describe its content in detail. If it contains text, tables, charts, diagrams, or any data, extract and present them clearly. Be thorough and precise."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_base64}"
                            }
                        }
                    ]
                }
            ],
            "max_completion_tokens": 4096
            # GPT-5.2 不支持 temperature=0，使用默认值 1
        }
        
        api_url = f"{AZURE_OPENAI_ENDPOINT.rstrip('/')}/openai/deployments/{AZURE_OPENAI_DEPLOYMENT}/chat/completions?api-version=2024-10-21"
        
        response = requests.post(
            api_url,
            headers=headers,
            json=payload,
            timeout=120
        )
        
        if response.status_code == 200:
            result = response.json()
            text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
            
            if len(text) > MAX_CONTENT_LENGTH:
                text = text[:MAX_CONTENT_LENGTH] + "\n... [内容截断]"
            
            return f"[图片分析 ({filename})]\n{text}" if text.strip() else "[图片无可分析内容]"
        else:
            error_msg = response.text[:200]
            return f"[GPT-5.2 视觉分析错误: HTTP {response.status_code} - {error_msg}]"
        
    except Exception as e:
        return f"[图片分析错误: {e}]"


def parse_zip(content: bytes, filename: str) -> str:
    """解析 ZIP 文件中的内容"""
    try:
        result_parts = []
        
        with zipfile.ZipFile(BytesIO(content)) as zf:
            for name in zf.namelist():
                ext = Path(name).suffix.lower()
                
                if ext in SKIP_EXTENSIONS:
                    result_parts.append(f"\n[ZIP 内跳过: {name} (音视频/特殊格式)]")
                    continue
                
                try:
                    file_content = zf.read(name)
                    
                    if ext == '.xlsx':
                        parsed = parse_excel(file_content, name)
                    elif ext == '.pdf':
                        parsed = parse_pdf(file_content, name)
                    elif ext == '.docx':
                        parsed = parse_docx(file_content, name)
                    elif ext == '.txt':
                        parsed = parse_txt(file_content, name)
                    elif ext in IMAGE_EXTENSIONS:
                        parsed = parse_image(file_content, name)
                    else:
                        parsed = f"[ZIP 内不支持的格式: {name}]"
                    
                    result_parts.append(f"\n=== ZIP/{name} ===\n{parsed}")
                    
                except Exception as e:
                    result_parts.append(f"\n[ZIP 内文件解析失败: {name}, {e}]")
        
        return '\n'.join(result_parts)
        
    except Exception as e:
        return f"[ZIP 解析错误: {e}]"


def parse_file(content: bytes, filename: str) -> str:
    """根据扩展名选择解析器"""
    ext = Path(filename).suffix.lower()
    
    if ext == '.xlsx':
        return parse_excel(content, filename)
    elif ext == '.pdf':
        return parse_pdf(content, filename)
    elif ext == '.docx':
        return parse_docx(content, filename)
    elif ext == '.pptx':
        return parse_pptx(content, filename)
    elif ext == '.txt':
        return parse_txt(content, filename)
    elif ext == '.zip':
        return parse_zip(content, filename)
    elif ext in IMAGE_EXTENSIONS:
        return parse_image(content, filename)
    elif ext in SKIP_EXTENSIONS:
        return f"[跳过: {filename} (音视频/特殊格式)]"
    else:
        return f"[不支持的文件格式: {ext}]"


# ============================================================
# 主处理逻辑
# ============================================================

def process_task(task: dict) -> dict:
    """处理单个任务，下载并解析附件"""
    enhanced_task = task.copy()
    
    urls = task.get('reference_file_urls', [])
    if not urls:
        enhanced_task['enhanced_prompt'] = task['prompt']
        enhanced_task['attachment_contents'] = []
        enhanced_task['attachment_summary'] = "无附件"
        return enhanced_task
    
    attachment_contents = []
    total_content = ""
    
    for url in urls:
        filename = url.split('/')[-1]
        ext = Path(filename).suffix.lower()
        
        # 下载
        content = download_file(url)
        if content is None:
            attachment_contents.append({
                'filename': filename,
                'status': 'download_failed',
                'content': None
            })
            continue
        
        # 解析
        parsed = parse_file(content, filename)
        
        attachment_contents.append({
            'filename': filename,
            'status': 'success' if not parsed.startswith('[') else 'partial',
            'content': parsed,
            'size_bytes': len(content),
            'parsed_length': len(parsed)
        })
        
        # 累积内容
        if len(total_content) < MAX_TOTAL_CONTENT:
            remaining = MAX_TOTAL_CONTENT - len(total_content)
            total_content += f"\n\n=== 附件: {filename} ===\n{parsed[:remaining]}"
    
    # 生成增强 prompt
    enhanced_prompt = task['prompt']
    if total_content:
        enhanced_prompt += f"\n\n{'='*60}\n以下是附件内容，请根据这些数据完成任务:\n{'='*60}\n{total_content}"
    
    enhanced_task['enhanced_prompt'] = enhanced_prompt
    enhanced_task['attachment_contents'] = attachment_contents
    enhanced_task['attachment_summary'] = f"{len(urls)} 个附件, {sum(1 for a in attachment_contents if a['status']=='success')} 个成功解析"
    
    return enhanced_task


def main():
    """主函数"""
    print("=" * 60)
    print("GDPVAL 数据增强")
    print(f"开始时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)
    
    # 加载原始数据
    input_path = Path(__file__).parent / "gdpval.json"
    with open(input_path, 'r', encoding='utf-8') as f:
        tasks = json.load(f)
    
    print(f"\n📊 总任务数: {len(tasks)}")
    print(f"📁 有附件: {sum(1 for t in tasks if t.get('reference_file_urls'))}")
    print(f"📂 缓存目录: {CACHE_DIR}")
    print()
    
    # 处理每个任务
    enhanced_tasks = []
    
    for i, task in enumerate(tasks):
        task_id = task['task_id'][:8]
        occupation = task['occupation'][:30]
        urls = task.get('reference_file_urls', [])
        
        print(f"\n[{i+1}/{len(tasks)}] {task_id}... {occupation}")
        
        if urls:
            print(f"  附件: {len(urls)} 个")
        
        enhanced = process_task(task)
        enhanced_tasks.append(enhanced)
        
        # 显示结果
        for att in enhanced.get('attachment_contents', []):
            status_icon = '✅' if att['status'] == 'success' else '⚠️'
            print(f"  {status_icon} {att['filename']}: {att.get('parsed_length', 0)} 字符")
    
    # 保存增强数据
    output_path = Path(__file__).parent / "gdpval_enhanced.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(enhanced_tasks, f, ensure_ascii=False, indent=2)
    
    print("\n" + "=" * 60)
    print("✅ 完成!")
    print(f"输出文件: {output_path}")
    print(f"结束时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    # 统计
    success_count = sum(
        1 for t in enhanced_tasks 
        for a in t.get('attachment_contents', []) 
        if a.get('status') == 'success'
    )
    total_attachments = sum(
        len(t.get('reference_file_urls', []))
        for t in tasks
    )
    
    print(f"\n📈 附件解析统计:")
    print(f"  总附件数: {total_attachments}")
    print(f"  成功解析: {success_count}")
    print(f"  成功率: {success_count/total_attachments*100:.1f}%" if total_attachments else "N/A")


if __name__ == "__main__":
    main()
