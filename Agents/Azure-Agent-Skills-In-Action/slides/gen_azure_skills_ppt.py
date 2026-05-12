"""Generate Azure Agent Skills In Action PPT — based on real 63-tool MCP run."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x1A, 0x1A, 0x3E)
GREEN = RGBColor(0x00, 0xB0, 0x4F)
ORANGE = RGBColor(0xFF, 0xAA, 0x00)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY = RGBColor(0x60, 0x60, 0x70)
LIGHT = RGBColor(0xF3, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return tb


def slide_header(slide, eyebrow, title):
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.05), BLUE)
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.35),
             eyebrow.upper(), size=11, bold=True, color=BLUE)
    add_text(slide, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
             title, size=28, bold=True, color=DARK)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Azure Agent Skills In Action  |  github.com/david-xinyuwei/david-share",
             size=9, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page} / 12", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# Slide 1: Cover
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
add_rect(s, 0, Inches(3.5), prs.slide_width, Inches(0.05), BLUE)
add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         "MICROSOFT AGENT SKILLS — INDEPENDENT EVALUATION",
         size=14, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.2),
         "Azure Agent Skills In Action",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
         "A 63-tool live run against a real Azure subscription",
         size=20, color=RGBColor(0xCC, 0xCC, 0xDD))
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
         "Xinyu Wei (魏新宇)  ·  Microsoft AI & Apps GBB",
         size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.4),
         "May 2026", size=12, color=RGBColor(0xAA, 0xAA, 0xCC))

# Slide 2: Why
s = prs.slides.add_slide(BLANK)
slide_header(s, "Context", "Why we built this evaluation")
add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         "Microsoft published two skills repos. Most teams will never run them end-to-end.",
         size=18, color=GRAY)
y = Inches(2.5)
items = [
    ("microsoft/azure-skills v1.1.39",
     "26 top-level skills · Azure MCP server · Foundry MCP — but no public proof of value."),
    ("microsoft/skills (174 skills)",
     "Python / .NET / TypeScript / Java / Rust SDK skills. M365 SDK skills are dev frameworks, not callable tools."),
    ("Our angle",
     "Run every callable Azure MCP tool against a real subscription. Publish raw evidence so others don't have to."),
]
for title, body in items:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.2), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.2), BLUE)
    add_text(s, Inches(0.85), y + Inches(0.15), Inches(11.5), Inches(0.4),
             title, size=16, bold=True, color=DARK)
    add_text(s, Inches(0.85), y + Inches(0.55), Inches(11.5), Inches(0.6),
             body, size=13, color=GRAY)
    y += Inches(1.4)
footer(s, 2)

# Slide 3: Test environment
s = prs.slides.add_slide(BLANK)
slide_header(s, "Test environment", "Real Azure subscription, Owner permission")
rows = [
    ("Subscription", "ME-MngEnv183724-xinyuwei-1"),
    ("Tenant", "MngEnv183724.onmicrosoft.com"),
    ("Permission", "Owner"),
    ("MCP server", "@azure/mcp@latest (npx)"),
    ("Protocol", "JSON-RPC 2024-11-05 over stdio"),
    ("Live resources", "8 VMs · 19 Cognitive Services · 20 Log Analytics · 10 Storage · 8 ML workspaces"),
    ("Auth", "Azure CLI token (current login)"),
    ("Run date", "2026-05-12"),
]
y = Inches(1.7)
for k, v in rows:
    add_rect(s, Inches(0.5), y, Inches(3.2), Inches(0.5), DARK)
    add_text(s, Inches(0.65), y + Inches(0.1), Inches(3), Inches(0.4),
             k, size=12, bold=True, color=WHITE)
    add_rect(s, Inches(3.7), y, Inches(9.1), Inches(0.5), LIGHT)
    add_text(s, Inches(3.85), y + Inches(0.1), Inches(8.9), Inches(0.4),
             v, size=12, color=DARK)
    y += Inches(0.6)
footer(s, 3)

# Slide 4: Headline
s = prs.slides.add_slide(BLANK)
slide_header(s, "Headline result", "63 / 63 tools probed — what actually executed")
stats = [
    ("45", "EXECUTED", GREEN, "Live Azure data returned"),
    ("9", "SCHEMA VERIFIED", BLUE, "Valid schema, missing resource input"),
    ("5", "TOOL ERROR", ORANGE, "Service / prerequisite issue"),
    ("2", "BLOCKED UNSAFE", RED, "Side-effecting, intentionally skipped"),
    ("2", "FAILED", GRAY, "Needs better test case"),
]
x = Inches(0.5); card_w = Inches(2.46); gap = Inches(0.05); y = Inches(1.7)
for num, label, color, desc in stats:
    add_rect(s, x, y, card_w, Inches(2.4), LIGHT)
    add_rect(s, x, y, card_w, Inches(0.15), color)
    add_text(s, x, y + Inches(0.45), card_w, Inches(1.2),
             num, size=64, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.55), card_w, Inches(0.4),
             label, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.95), card_w - Inches(0.2), Inches(0.5),
             desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    x += card_w + gap
add_rect(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.0), DARK)
add_text(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.5),
         "Coverage interpretation", size=16, bold=True, color=BLUE)
add_text(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.3),
         "63 / 63 top-level tools probed.   45 executed live.   54 produced live execution evidence "
         "or a verified command schema (45 + 9).   The remaining 9 are documented with the exact blocker.",
         size=14, color=WHITE)
footer(s, 4)

# Slide 5: Wins
s = prs.slides.add_slide(BLANK)
slide_header(s, "What worked", "High-signal wins from the run")
wins = [
    ("Live inventory", "subscription_list · group_list · group_resource_list",
     "Real Azure state through current CLI login."),
    ("Compute & apps", "compute_vm_get · aks_cluster_get · containerapps_list · appservice_webapp_get",
     "Inspect runtime infra without hand-writing az queries."),
    ("Cost / quota / pricing", "quota_usage_check · pricing_get · advisor_recommendation_list",
     "High-friction APIs delivered through one consistent contract."),
    ("IaC & architecture", "bicepschema_get · azureterraform_azurerm_get · cloudarchitect_design",
     "Schemas and architecture guidance on demand."),
    ("Governance / identity", "role_assignment_list · policy_assignment_list · resourcehealth_*",
     "RBAC, policy, health surfaced as structured evidence."),
    ("Service discovery", "storage_account_get · cosmos_list · sql_server_get · redis_list · search_service_list",
     "Same contract sweeps many service families."),
]
y = Inches(1.6)
for title, tools, desc in wins:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.12), Inches(0.85), GREEN)
    add_text(s, Inches(0.85), y + Inches(0.05), Inches(3.0), Inches(0.4),
             title, size=12, bold=True, color=DARK)
    add_text(s, Inches(0.85), y + Inches(0.42), Inches(11.5), Inches(0.4),
             tools, size=10, color=BLUE)
    add_text(s, Inches(3.95), y + Inches(0.05), Inches(8.6), Inches(0.4),
             desc, size=11, color=GRAY)
    y += Inches(0.92)
footer(s, 5)

# Slide 6: Didn't fully execute
s = prs.slides.add_slide(BLANK)
slide_header(s, "What didn't fully execute", "And exactly why — no hand-waving")
rows = [
    ("Needs resource instance", BLUE,
     "keyvault · servicebus · servicefabric · speech · foundryextensions · confidentialledger · datadog · mysql · deploy",
     "Schema is valid. Need a vault, queue, speech file, endpoint, cluster, ledger TX, Datadog resource, MySQL user, or local azd workspace."),
    ("Intentionally skipped", RED,
     "communication · azuremigrate",
     "Selected commands could send SMS or trigger migration changes. Schema recorded; execution blocked by design."),
    ("Product / prerequisite", ORANGE,
     "extension_azqr · loadtesting · marketplace · applens · foundry",
     "Server returned runtime error or prerequisite missing. Example: extension_azqr requires the azqr binary in PATH."),
    ("Needs better test case", GRAY,
     "applicationinsights · extension_cli_install",
     "Argument set we tried wasn't sufficient. Listed in matrix so others can improve it."),
]
y = Inches(1.6)
for cat, color, tools, why in rows:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.2), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(1.2), color)
    add_text(s, Inches(0.85), y + Inches(0.1), Inches(11.5), Inches(0.4),
             cat, size=14, bold=True, color=color)
    add_text(s, Inches(0.85), y + Inches(0.45), Inches(11.5), Inches(0.4),
             tools, size=11, color=DARK)
    add_text(s, Inches(0.85), y + Inches(0.78), Inches(11.5), Inches(0.4),
             why, size=10, color=GRAY)
    y += Inches(1.3)
footer(s, 6)

# Slide 7: Calling convention
s = prs.slides.add_slide(BLANK)
slide_header(s, "Calling convention", "Direct JSON-RPC: composite tools take flat args")
add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.6),
         "What we proved by running it: composite tools accept command + flat arguments,"
         " not a JSON-string parameters wrapper.",
         size=14, color=GRAY)
add_rect(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.4), DARK)
code = ("send(\"compute\", {\n"
        "  command: \"compute_vm_get\",\n"
        "  subscription: SUB,\n"
        "  \"resource-group\": \"winvm\"\n"
        "});")
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(code.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.name = "Consolas"; r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(0xC8, 0xE8, 0xFF)
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), LIGHT)
add_text(s, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.4),
         "Why this matters", size=14, bold=True, color=DARK)
add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0),
         "The mcp_azure_mcp_* prefix in SKILL.md is added by the host (VS Code, Copilot CLI). "
         "The raw server exposes plain names: compute, quota, pricing, subscription_list, group_list. "
         "We corrected an earlier wrong assumption (JSON-string parameters wrapper) by actually running it.",
         size=12, color=GRAY)
footer(s, 7)

# Slide 8: Skills vs CLI
s = prs.slides.add_slide(BLANK)
slide_header(s, "Skills vs az CLI", "Two concrete cases from the run")
add_rect(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.5), LIGHT)
add_text(s, Inches(0.65), Inches(1.7), Inches(5.7), Inches(0.4),
         "Case 1 — list subscriptions", size=14, bold=True, color=DARK)
add_text(s, Inches(0.65), Inches(2.1), Inches(5.7), Inches(0.4),
         "Without skills (az CLI):", size=11, bold=True, color=BLUE)
add_text(s, Inches(0.65), Inches(2.45), Inches(5.7), Inches(0.4),
         "az account list — 0.95s, JSON output", size=11, color=GRAY)
add_text(s, Inches(0.65), Inches(2.9), Inches(5.7), Inches(0.4),
         "With skills (subscription_list):", size=11, bold=True, color=GREEN)
add_text(s, Inches(0.65), Inches(3.25), Inches(5.7), Inches(0.4),
         "8 ms, structured JSON ready for LLM", size=11, color=GRAY)
add_text(s, Inches(0.65), Inches(3.7), Inches(5.7), Inches(0.4),
         "Verdict: similar speed, MCP wins for LLM consumption.",
         size=11, bold=True, color=DARK)

add_rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(2.5), LIGHT)
add_text(s, Inches(6.95), Inches(1.7), Inches(5.7), Inches(0.4),
         "Case 2 — quota usage", size=14, bold=True, color=DARK)
add_text(s, Inches(6.95), Inches(2.1), Inches(5.7), Inches(0.4),
         "Without skills:", size=11, bold=True, color=BLUE)
add_text(s, Inches(6.95), Inches(2.45), Inches(5.7), Inches(0.4),
         "Manual REST + signed token + custom parsing",
         size=11, color=GRAY)
add_text(s, Inches(6.95), Inches(2.9), Inches(5.7), Inches(0.4),
         "With skills (quota_usage_check):", size=11, bold=True, color=GREEN)
add_text(s, Inches(6.95), Inches(3.25), Inches(5.7), Inches(0.4),
         "One call → 18.6 KB structured quota report",
         size=11, color=GRAY)
add_text(s, Inches(6.95), Inches(3.7), Inches(5.7), Inches(0.4),
         "Verdict: skill wins by an order of magnitude.",
         size=11, bold=True, color=DARK)

add_rect(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.0), DARK)
add_text(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.4),
         "When skills win", size=14, bold=True, color=BLUE)
add_bullets(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(1.4),
            [
                "High-friction APIs (quota, pricing, advisor) wrapped in one consistent contract.",
                "Cross-service architecture / Bicep / Terraform schemas on demand.",
                "Guardrails: read-only first, ask before destructive — built into the harness.",
            ], size=12, color=WHITE)
footer(s, 8)

# Slide 9: Architecture
s = prs.slides.add_slide(BLANK)
slide_header(s, "Architecture", "How a single Azure MCP server fronts everything")

def block(x, y, w, h, title, sub, color):
    add_rect(s, x, y, w, h, color)
    add_text(s, x, y + Inches(0.1), w, Inches(0.4),
             title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.45), w, Inches(0.4),
             sub, size=10, color=RGBColor(0xE0, 0xE8, 0xFF), align=PP_ALIGN.CENTER)

block(Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.85),
      "Agent Host (VS Code Copilot · Claude Code · Custom Agent)",
      "Adds mcp_azure_mcp_* prefix · handles auth handoff", BLUE)
add_text(s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.3),
         "▼  JSON-RPC 2024-11-05 over stdio  ▼",
         size=11, color=GRAY, align=PP_ALIGN.CENTER)
block(Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.85),
      "@azure/mcp server  (npx -y @azure/mcp@latest server start)",
      "63 top-level tools · simple + composite (learn → execute) · read-only safe", DARK)
add_text(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.3),
         "▼  Azure CLI / SDK auth  ▼",
         size=11, color=GRAY, align=PP_ALIGN.CENTER)
fams = [
    ("Compute & Apps", BLUE), ("Data & Storage", GREEN),
    ("Ops & Governance", ORANGE), ("Identity & Sec.", RED),
    ("AI & Foundry", BLUE), ("IaC / Architect", GREEN),
]
fx = Inches(0.5); fw = Inches(2.0); fgap = Inches(0.06)
for name, color in fams:
    add_rect(s, fx, Inches(4.3), fw, Inches(0.7), color)
    add_text(s, fx, Inches(4.45), fw, Inches(0.4),
             name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    fx += fw + fgap
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.4), LIGHT)
add_text(s, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.4),
         "Key insight", size=14, bold=True, color=DARK)
add_text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.8),
         "One MCP server fronts every Azure service family. The friction we observed is at the "
         "input-parameter layer — host-injected wrappers (mcp_azure_mcp_* + auto args) hide it from end users.",
         size=12, color=GRAY)
footer(s, 9)

# Slide 10: Stickiness
s = prs.slides.add_slide(BLANK)
slide_header(s, "Platform stickiness", "What locks you into Microsoft once you adopt these")
items = [
    ("HIGH",  RED,    "cloudarchitect_design", "Reasoning about Azure-only services; port = rebuild."),
    ("HIGH",  RED,    "wellarchitectedframework", "WAF guidance is Azure-specific by design."),
    ("HIGH",  RED,    "bicepschema · azureterraform", "Generated IaC targets azurerm provider."),
    ("MED",   ORANGE, "advisor · quota · pricing", "Same surface in other clouds, but Azure-shaped."),
    ("MED",   ORANGE, "policy · role", "RBAC + Policy semantics are Azure RM only."),
    ("LOW",   GREEN,  "subscription_list · group_list · *_list", "Thin wrappers; easy to swap."),
]
y = Inches(1.6)
for level, color, tools, desc in items:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.75), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(1.0), Inches(0.75), color)
    add_text(s, Inches(0.5), y + Inches(0.18), Inches(1.0), Inches(0.4),
             level, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.65), y + Inches(0.08), Inches(11.0), Inches(0.4),
             tools, size=12, bold=True, color=DARK)
    add_text(s, Inches(1.65), y + Inches(0.4), Inches(11.0), Inches(0.4),
             desc, size=10, color=GRAY)
    y += Inches(0.85)
footer(s, 10)

# Slide 11: Deliverables
s = prs.slides.add_slide(BLANK)
slide_header(s, "What we deliver", "The repo other engineers can clone today")
cards = [
    ("scripts/run_full_value_evaluation.js",
     "Full 63-tool harness · learn → safe execute · classified outputs.", BLUE),
    ("evaluation/results/full_value_evaluation.json",
     "725-line raw JSON evidence — every tool, every input, every output.", GREEN),
    ("evaluation/results/full_value_matrix.csv",
     "63-row matrix: tool · family · mode · status · duration · output size.", ORANGE),
    ("evaluation/results/full_value_summary.md",
     "Human-readable markdown report linked from README.", BLUE),
    ("README.md / README-CN.md",
     "653 lines each, line-aligned, full-run evidence as the central narrative.", DARK),
]
y = Inches(1.6)
for path, desc, color in cards:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.95), color)
    add_text(s, Inches(0.85), y + Inches(0.1), Inches(11.5), Inches(0.4),
             path, size=13, bold=True, color=DARK)
    add_text(s, Inches(0.85), y + Inches(0.5), Inches(11.5), Inches(0.4),
             desc, size=11, color=GRAY)
    y += Inches(1.05)
footer(s, 11)

# Slide 12: Closing
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
add_rect(s, 0, Inches(2.3), prs.slide_width, Inches(0.05), BLUE)
add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
         "TAKE AWAY", size=14, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.7),
         "Don't trust skills catalogs. Run them.",
         size=36, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.5),
         "What this evaluation gives you",
         size=16, bold=True, color=BLUE)
add_bullets(s, Inches(0.7), Inches(3.15), Inches(12), Inches(2.4),
            [
                "Live evidence that 45 / 63 Azure MCP tools work end-to-end against a real subscription.",
                "Honest disclosure of 9 schema-only / 5 tool-error / 2 blocked / 2 failed cases.",
                "Corrected calling convention for composite tools — verified by execution.",
                "Reusable harness so others can re-run against their own subscriptions.",
            ], size=14, color=WHITE)
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "Repo:  github.com/david-xinyuwei/david-share/tree/master/Agents/Azure-Agent-Skills-In-Action",
         size=12, color=RGBColor(0xCC, 0xE8, 0xFF))
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
         "Author:  Xinyu Wei (魏新宇)  ·  Microsoft AI & Apps GBB  ·  May 2026",
         size=12, color=RGBColor(0xAA, 0xAA, 0xCC))

import os
out = "/mnt/g/AI-Super-Agent/slides-generate/working/Azure-Agent-Skills-In-Action.pptx"
prs.save(out)
print(f"OK saved: {out}  size={os.path.getsize(out)} bytes  slides={len(prs.slides)}")
